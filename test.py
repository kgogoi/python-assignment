from pptx.util import Inches, Pt
from wand.compat import text
from wand.image import Image
from pptx import Presentation
from os import chdir, listdir
import pytest;

# Counts number of images to be itirated over
img_count = len(listdir('./input_img'));

# pptx init
prs = Presentation();
slide_register = prs.slide_layouts[6];

# Resizes image to smaller dimensions adds adds logo
def img_transform(i):
    with Image(filename = './input_img/image{}.jpg'.format(i+1)) as img1:
        with Image(filename = 'nike_black.png') as img2:
            img1.transform(resize = '5.5%');
            img2.transform(resize = '5%');
            img1.composite(image=img2, left=0, top=0)
            img1.save(filename = './output_img/image-edited{}.jpg'.format(i+1));

# Iterates over each slide and adds logo and heading to it
def slide_create(prs, slide_register, i):
    slide = prs.slides.add_slide(slide_register);
    # Adding Image to each slide
    left = Inches(1);
    top = Inches(1.65);
    slide.shapes.add_picture('./output_img/image-edited{}.jpg'.format(i+1), left, top);
    left = width = height = Inches(.6);
    top = Inches(.3);
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame;
    # Add Heading
    p = tf.add_paragraph();
    p.text = "Sample Heading {}".format(i+1);
    p.font.size = Pt(28)
    # Add Subheading
    p = tf.add_paragraph();
    p.text = "Sample Subheading {}".format(i+1);
    p.font.size = Pt(16)

for i in range(img_count):
    img_transform(i);
    slide_create(prs, slide_register, i)
    prs.save('assignment.pptx');
